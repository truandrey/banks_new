from pptx import Presentation
from pptx.util import Inches
from datetime import date

def createPresentation():
    # Инициализация презентации формата pptx
    ppt = Presentation()

    # Создание титульного слайда
    first_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    title = "Статистический отчет - " + str(date.today())
    first_slide.shapes[0].text_frame.paragraphs[0].text = title

    # Создание слайда со статистикой
    img = 'sources/BarChart.png'
    second_slide = ppt.slide_layouts[1]
    slide2 = ppt.slides.add_slide(second_slide)
    slide2.shapes.add_picture(img, left = Inches(0.3), top = Inches(0.2), height = Inches(7))
    ppt.save('reviews/Report.pptx')